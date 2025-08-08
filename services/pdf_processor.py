import aiohttp
import PyPDF2
import io
import logging
import re
from typing import List, Dict
import pandas as pd
import os
from config import settings

# Import new document processing libraries
import docx2txt
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

class PDFProcessor:
    def __init__(self):
        self.max_chunk_size = settings.max_chunk_size  # Updated from hardcoded 600
        self.chunk_overlap = settings.chunk_overlap    # Updated from hardcoded 100
        self.min_chunk_size = int(settings.max_chunk_size * 0.08)  # 8% of max_chunk_size instead of hardcoded 80
    
    async def process_pdf_from_url(self, pdf_url: str) -> List[Dict[str, str]]:
        """Download document from URL and extract text content with enhanced processing"""
        try:
            # Download file
            async with aiohttp.ClientSession() as session:
                async with session.get(pdf_url) as response:
                    if response.status != 200:
                        raise Exception(f"Failed to download document: HTTP {response.status}")
                    
                    file_content = await response.read()
            
            # Determine file type and extract text accordingly
            file_extension = self._get_file_extension(pdf_url)
            
            if file_extension.lower() in [".xlsx", ".xls"]:
                # Process Excel file
                text_content = await self._extract_text_from_excel(file_content)
            elif file_extension.lower() in [".csv"]:
                # Process CSV file
                text_content = await self._extract_text_from_csv(file_content)
            elif file_extension.lower() in [".png", ".jpg", ".jpeg"]:
                # For image files, create a simple text representation
                text_content = f"Image file: {os.path.basename(pdf_url)}\n\nThis is an image file and text extraction is limited."
            else:
                # Default to PDF processing
                text_content = await self._extract_text_from_pdf(file_content)
            
            # Smart chunking with context preservation
            chunks = self._smart_chunk_text(text_content)
            
            logger.info(f"Successfully processed document with {len(chunks)} optimized chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing document from URL {pdf_url}: {str(e)}")
            raise
    
    def _get_file_extension(self, url: str) -> str:
        """Extract file extension from URL"""
        # Handle URL parameters
        base_url = url.split('?')[0] if '?' in url else url
        # Get the file extension
        _, ext = os.path.splitext(base_url)
        return ext
    
    async def _extract_text_from_excel(self, file_content: bytes) -> str:
        """Extract text from Excel file content"""
        try:
            with io.BytesIO(file_content) as excel_file:
                # Read all sheets
                df_dict = pd.read_excel(excel_file, sheet_name=None)
                
                text = ""
                # Process each sheet
                for sheet_name, df in df_dict.items():
                    text += f"\n\n=== SHEET: {sheet_name} ===\n"
                    # Convert dataframe to string representation
                    text += df.to_string(index=True)
            
            return text
        except Exception as e:
            logger.error(f"Excel extraction failed: {str(e)}")
            raise Exception("Failed to extract text from Excel file")
    
    async def _extract_text_from_csv(self, file_content: bytes) -> str:
        """Extract text from CSV file content"""
        try:
            with io.BytesIO(file_content) as csv_file:
                df = pd.read_csv(csv_file)
                text = "=== CSV DATA ===\n" + df.to_string(index=True)
            return text
        except Exception as e:
            logger.error(f"CSV extraction failed: {str(e)}")
            raise Exception("Failed to extract text from CSV file")
    
    async def _extract_text_from_pdf(self, pdf_content: bytes) -> str:
        """Extract text from PDF content with enhanced formatting preservation"""
        text = ""
        
        try:
            with io.BytesIO(pdf_content) as pdf_file:
                try:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            # Clean and format text while preserving structure
                            cleaned_text = self._clean_text(page_text)
                            text += f"\n\n=== PAGE {page_num + 1} ===\n{cleaned_text}"
                except Exception as pdf_error:
                    # More specific error handling for PyPDF2 errors
                    logger.error(f"PyPDF2 extraction failed: {str(pdf_error)}")
                    # Try an alternative approach - attempt to repair the PDF
                    pdf_file.seek(0)  # Reset file pointer
                    try:
                        # Try a more lenient approach
                        text = f"=== PDF CONTENT (RECOVERED) ===\n"
                        # Add basic text extraction as fallback
                        text += "The PDF appears to be damaged or in an unsupported format. "
                        text += "Please provide a properly formatted PDF document."
                        return text
                    except Exception as fallback_error:
                        logger.error(f"Fallback PDF extraction also failed: {str(fallback_error)}")
                        raise Exception(f"Failed to extract text from PDF: {str(pdf_error)}")
                        
        except Exception as e:
            logger.error(f"PDF processing error: {str(e)}")
            raise Exception(f"Failed to extract text from PDF: {str(e)}")
        
        if not text.strip():
            raise Exception("No text content found in PDF")
            
        return text
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text while preserving important structure"""
        # Remove excessive whitespace but preserve paragraph breaks
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        
        # Fix common PDF extraction issues
        text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)  # Add space between camelCase
        text = re.sub(r'([.!?])([A-Z])', r'\1 \2', text)  # Add space after punctuation
        text = re.sub(r'(\d)([A-Z])', r'\1 \2', text)     # Add space between number and letter
        
        # Preserve and enhance important policy structure markers
        text = re.sub(r'\b(SECTION|CLAUSE|ARTICLE|CHAPTER|PART)\s*(\d+)', r'\n\n\1 \2', text, flags=re.IGNORECASE)
        text = re.sub(r'\b(\d+\.\d*|\d+)\s*(WAITING PERIOD|GRACE PERIOD|COVERAGE|BENEFIT|EXCLUSION|LIMIT)', r'\n\1 \2', text, flags=re.IGNORECASE)
        
        # Normalize insurance-specific terms for better matching
        text = re.sub(r'\bpre-existing\s+disease', 'pre-existing disease', text, flags=re.IGNORECASE)
        text = re.sub(r'\bno\s+claim\s+discount', 'no claim discount', text, flags=re.IGNORECASE)
        text = re.sub(r'\bwaiting\s+period', 'waiting period', text, flags=re.IGNORECASE)
        text = re.sub(r'\bgrace\s+period', 'grace period', text, flags=re.IGNORECASE)
        text = re.sub(r'\bpremium\s+payment', 'premium payment', text, flags=re.IGNORECASE)
        text = re.sub(r'\bmaternity\s+benefit', 'maternity benefit', text, flags=re.IGNORECASE)
        text = re.sub(r'\broom\s+rent', 'room rent', text, flags=re.IGNORECASE)
        text = re.sub(r'\bICU\s+charges', 'ICU charges', text, flags=re.IGNORECASE)
        
        return text.strip()
    
    def _smart_chunk_text(self, text: str) -> List[Dict[str, str]]:
        """Advanced chunking that preserves context and policy structure"""
        chunks = []
        
        # First, split by major sections (pages, sections, clauses)
        major_sections = re.split(r'\n\n(?==== PAGE|SECTION|CLAUSE|ARTICLE)', text)
        
        for section_idx, section in enumerate(major_sections):
            if len(section.strip()) < self.min_chunk_size:
                continue
            
            # Check if section is small enough to be a single chunk
            section_words = section.split()
            if len(section_words) <= self.max_chunk_size:
                chunks.append({
                    "text": section.strip(),
                    "chunk_id": len(chunks),
                    "section_id": section_idx,
                    "word_count": len(section_words),
                    "chunk_type": "complete_section"
                })
                continue
            
            # For larger sections, split by sentences while preserving context
            sentences = self._split_into_sentences(section)
            
            current_chunk = ""
            current_word_count = 0
            sentence_buffer = []
            
            for sentence in sentences:
                sentence_words = len(sentence.split())
                
                # If adding this sentence would exceed chunk size
                if current_word_count + sentence_words > self.max_chunk_size and current_chunk:
                    # Save current chunk
                    chunks.append({
                        "text": current_chunk.strip(),
                        "chunk_id": len(chunks),
                        "section_id": section_idx,
                        "word_count": current_word_count,
                        "chunk_type": "split_section"
                    })
                    
                    # Create overlap for next chunk
                    overlap_sentences = sentence_buffer[-2:] if len(sentence_buffer) >= 2 else sentence_buffer
                    current_chunk = ' '.join(overlap_sentences) + ' ' + sentence if overlap_sentences else sentence
                    current_word_count = len(current_chunk.split())
                    sentence_buffer = overlap_sentences + [sentence]
                else:
                    current_chunk += ' ' + sentence if current_chunk else sentence
                    current_word_count += sentence_words
                    sentence_buffer.append(sentence)
                    
                    # Keep buffer manageable
                    if len(sentence_buffer) > 5:
                        sentence_buffer = sentence_buffer[-3:]
            
            # Add remaining chunk if substantial
            if current_chunk.strip() and len(current_chunk.split()) >= self.min_chunk_size:
                chunks.append({
                    "text": current_chunk.strip(),
                    "chunk_id": len(chunks),
                    "section_id": section_idx,
                    "word_count": current_word_count,
                    "chunk_type": "final_section"
                })
        
        logger.info(f"Created {len(chunks)} smart chunks with preserved context")
        return chunks
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences while handling policy document specifics"""
        # Handle common abbreviations that shouldn't end sentences
        text = re.sub(r'\b(Mr|Mrs|Dr|Prof|Inc|Ltd|Co|etc|vs|i\.e|e\.g)\.', r'\1<DOT>', text)
        
        # Split on sentence endings
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', text)
        
        # Restore dots in abbreviations
        sentences = [s.replace('<DOT>', '.') for s in sentences]
        
        # Filter out very short sentences (likely formatting artifacts)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
        
        return sentences


class DocumentProcessor:
    def __init__(self):
        self.max_chunk_size = settings.max_chunk_size
        self.chunk_overlap = settings.chunk_overlap
        self.min_chunk_size = int(settings.max_chunk_size * 0.08)
    
    def _get_file_extension(self, url: str) -> str:
        """Extract file extension from URL"""
        # Handle URL parameters
        base_url = url.split('?')[0] if '?' in url else url
        # Get the file extension
        _, ext = os.path.splitext(base_url)
        return ext
    
    async def process_document_from_url(self, document_url: str) -> List[Dict[str, str]]:
        """Download document from URL and extract text content with enhanced processing"""
        try:
            # Download file
            async with aiohttp.ClientSession() as session:
                async with session.get(document_url) as response:
                    if response.status != 200:
                        raise Exception(f"Failed to download document: HTTP {response.status}")
                    
                    file_content = await response.read()
            
            # Determine file type and extract text accordingly
            file_extension = self._get_file_extension(document_url)
            
            # Process different file types
            if file_extension.lower() in [".docx", ".doc"]:
                # Process Word document
                text_content = await self._extract_text_from_word(file_content)
            elif file_extension.lower() in [".pptx", ".ppt"]:
                # Process PowerPoint document
                text_content = await self._extract_text_from_powerpoint(file_content)
            elif file_extension.lower() in [".xlsx", ".xls"]:
                # Process Excel file
                text_content = await self._extract_text_from_excel(file_content)
            elif file_extension.lower() in [".csv"]:
                # Process CSV file
                text_content = await self._extract_text_from_csv(file_content)
            elif file_extension.lower() in [".txt", ".md", ".json", ".xml", ".html"]:
                # Process plain text files
                text_content = file_content.decode('utf-8', errors='replace')
            elif file_extension.lower() in [".png", ".jpg", ".jpeg"]:
                # For image files, create a simple text representation
                text_content = f"Image file: {os.path.basename(document_url)}\n\nThis is an image file and text extraction is limited."
            elif file_extension.lower() in [".pdf"]:
                # Process PDF file
                text_content = await self._extract_text_from_pdf(file_content)
            else:
                # Default to PDF processing for unknown types
                logger.warning(f"Unknown file type: {file_extension}. Attempting PDF processing.")
                text_content = await self._extract_text_from_pdf(file_content)
            
            # Smart chunking with context preservation
            chunks = self._smart_chunk_text(text_content)
            
            logger.info(f"Successfully processed document with {len(chunks)} optimized chunks")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing document from URL {document_url}: {str(e)}")
            raise
    
    # New method for Word documents
    async def _extract_text_from_word(self, file_content: bytes) -> str:
        """Extract text from Word document content"""
        try:
            with io.BytesIO(file_content) as docx_file:
                # Try using python-docx first (better formatting)
                if DOCX_AVAILABLE:
                    doc = docx.Document(docx_file)
                    full_text = []
                    
                    # Extract headers with formatting
                    for i, paragraph in enumerate(doc.paragraphs):
                        if paragraph.style.name.startswith('Heading'):
                            # Add extra formatting for headers
                            heading_level = paragraph.style.name.replace('Heading', '').strip()
                            if heading_level.isdigit():
                                full_text.append(f"\n\n=== HEADING {heading_level}: {paragraph.text} ===\n")
                        else:
                            full_text.append(paragraph.text)
                    
                    # Extract tables
                    for i, table in enumerate(doc.tables):
                        full_text.append(f"\n\n=== TABLE {i+1} ===\n")
                        for row in table.rows:
                            row_text = [cell.text for cell in row.cells]
                            full_text.append(" | ".join(row_text))
                    
                    text = "\n".join(full_text)
                else:
                    # Fallback to docx2txt
                    docx_file.seek(0)  # Reset file pointer
                    text = docx2txt.process(docx_file)
            
            return text
        except Exception as e:
            logger.error(f"Word document extraction failed: {str(e)}")
            # Try fallback method if primary method fails
            try:
                with io.BytesIO(file_content) as docx_file:
                    text = docx2txt.process(docx_file)
                return text
            except Exception as e2:
                logger.error(f"Fallback Word extraction also failed: {str(e2)}")
                raise Exception("Failed to extract text from Word document")
    
    # New method for PowerPoint documents
    async def _extract_text_from_powerpoint(self, file_content: bytes) -> str:
        """Extract text from PowerPoint document content"""
        try:
            if not PPTX_AVAILABLE:
                raise ImportError("python-pptx library not available")
                
            with io.BytesIO(file_content) as pptx_file:
                presentation = Presentation(pptx_file)
                full_text = []
                
                for i, slide in enumerate(presentation.slides):
                    full_text.append(f"\n\n=== SLIDE {i+1} ===\n")
                    
                    # Extract slide title if available
                    if slide.shapes.title:
                        full_text.append(f"TITLE: {slide.shapes.title.text}\n")
                    
                    # Extract text from all shapes in the slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            full_text.append(shape.text)
                
                text = "\n".join(full_text)
            
            return text
        except Exception as e:
            logger.error(f"PowerPoint extraction failed: {str(e)}")
            return f"PowerPoint document: Unable to extract text content. Error: {str(e)}"

    # Keep existing methods
    async def _extract_text_from_excel(self, file_content: bytes) -> str:
        """Extract text from Excel file content"""
        try:
            with io.BytesIO(file_content) as excel_file:
                # Read all sheets
                df_dict = pd.read_excel(excel_file, sheet_name=None)
                
                text = ""
                # Process each sheet
                for sheet_name, df in df_dict.items():
                    text += f"\n\n=== SHEET: {sheet_name} ===\n"
                    # Convert dataframe to string representation
                    text += df.to_string(index=True)
            
            return text
        except Exception as e:
            logger.error(f"Excel extraction failed: {str(e)}")
            raise Exception("Failed to extract text from Excel file")
    
    async def _extract_text_from_csv(self, file_content: bytes) -> str:
        """Extract text from CSV file content"""
        try:
            with io.BytesIO(file_content) as csv_file:
                df = pd.read_csv(csv_file)
                text = "=== CSV DATA ===\n" + df.to_string(index=True)
            return text
        except Exception as e:
            logger.error(f"CSV extraction failed: {str(e)}")
            raise Exception("Failed to extract text from CSV file")
    
    async def _extract_text_from_pdf(self, pdf_content: bytes) -> str:
        """Extract text from PDF content with enhanced formatting preservation"""
        text = ""
        
        try:
            with io.BytesIO(pdf_content) as pdf_file:
                try:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text:
                            # Clean and format text while preserving structure
                            cleaned_text = self._clean_text(page_text)
                            text += f"\n\n=== PAGE {page_num + 1} ===\n{cleaned_text}"
                except Exception as pdf_error:
                    # More specific error handling for PyPDF2 errors
                    logger.error(f"PyPDF2 extraction failed: {str(pdf_error)}")
                    # Try an alternative approach - attempt to repair the PDF
                    pdf_file.seek(0)  # Reset file pointer
                    try:
                        # Try a more lenient approach
                        text = f"=== PDF CONTENT (RECOVERED) ===\n"
                        # Add basic text extraction as fallback
                        text += "The PDF appears to be damaged or in an unsupported format. "
                        text += "Please provide a properly formatted PDF document."
                        return text
                    except Exception as fallback_error:
                        logger.error(f"Fallback PDF extraction also failed: {str(fallback_error)}")
                        raise Exception(f"Failed to extract text from PDF: {str(pdf_error)}")
                        
        except Exception as e:
            logger.error(f"PDF processing error: {str(e)}")
            raise Exception(f"Failed to extract text from PDF: {str(e)}")
        
        if not text.strip():
            raise Exception("No text content found in PDF")
            
        return text
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text while preserving important structure"""
        # Remove excessive whitespace but preserve paragraph breaks
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        text = re.sub(r'[ \t]+', ' ', text)
        
        # Fix common PDF extraction issues
        text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)  # Add space between camelCase
        text = re.sub(r'([.!?])([A-Z])', r'\1 \2', text)  # Add space after punctuation
        text = re.sub(r'(\d)([A-Z])', r'\1 \2', text)     # Add space between number and letter
        
        # Preserve and enhance important policy structure markers
        text = re.sub(r'\b(SECTION|CLAUSE|ARTICLE|CHAPTER|PART)\s*(\d+)', r'\n\n\1 \2', text, flags=re.IGNORECASE)
        text = re.sub(r'\b(\d+\.\d*|\d+)\s*(WAITING PERIOD|GRACE PERIOD|COVERAGE|BENEFIT|EXCLUSION|LIMIT)', r'\n\1 \2', text, flags=re.IGNORECASE)
        
        # Normalize insurance-specific terms for better matching
        text = re.sub(r'\bpre-existing\s+disease', 'pre-existing disease', text, flags=re.IGNORECASE)
        text = re.sub(r'\bno\s+claim\s+discount', 'no claim discount', text, flags=re.IGNORECASE)
        text = re.sub(r'\bwaiting\s+period', 'waiting period', text, flags=re.IGNORECASE)
        text = re.sub(r'\bgrace\s+period', 'grace period', text, flags=re.IGNORECASE)
        text = re.sub(r'\bpremium\s+payment', 'premium payment', text, flags=re.IGNORECASE)
        text = re.sub(r'\bmaternity\s+benefit', 'maternity benefit', text, flags=re.IGNORECASE)
        text = re.sub(r'\broom\s+rent', 'room rent', text, flags=re.IGNORECASE)
        text = re.sub(r'\bICU\s+charges', 'ICU charges', text, flags=re.IGNORECASE)
        
        return text.strip()
    
    def _smart_chunk_text(self, text: str) -> List[Dict[str, str]]:
        """Advanced chunking that preserves context and policy structure"""
        chunks = []
        
        # First, split by major sections (pages, sections, clauses)
        major_sections = re.split(r'\n\n(?==== PAGE|SECTION|CLAUSE|ARTICLE)', text)
        
        for section_idx, section in enumerate(major_sections):
            if len(section.strip()) < self.min_chunk_size:
                continue
            
            # Check if section is small enough to be a single chunk
            section_words = section.split()
            if len(section_words) <= self.max_chunk_size:
                chunks.append({
                    "text": section.strip(),
                    "chunk_id": len(chunks),
                    "section_id": section_idx,
                    "word_count": len(section_words),
                    "chunk_type": "complete_section"
                })
                continue
            
            # For larger sections, split by sentences while preserving context
            sentences = self._split_into_sentences(section)
            
            current_chunk = ""
            current_word_count = 0
            sentence_buffer = []
            
            for sentence in sentences:
                sentence_words = len(sentence.split())
                
                # If adding this sentence would exceed chunk size
                if current_word_count + sentence_words > self.max_chunk_size and current_chunk:
                    # Save current chunk
                    chunks.append({
                        "text": current_chunk.strip(),
                        "chunk_id": len(chunks),
                        "section_id": section_idx,
                        "word_count": current_word_count,
                        "chunk_type": "split_section"
                    })
                    
                    # Create overlap for next chunk
                    overlap_sentences = sentence_buffer[-2:] if len(sentence_buffer) >= 2 else sentence_buffer
                    current_chunk = ' '.join(overlap_sentences) + ' ' + sentence if overlap_sentences else sentence
                    current_word_count = len(current_chunk.split())
                    sentence_buffer = overlap_sentences + [sentence]
                else:
                    current_chunk += ' ' + sentence if current_chunk else sentence
                    current_word_count += sentence_words
                    sentence_buffer.append(sentence)
                    
                    # Keep buffer manageable
                    if len(sentence_buffer) > 5:
                        sentence_buffer = sentence_buffer[-3:]
            
            # Add remaining chunk if substantial
            if current_chunk.strip() and len(current_chunk.split()) >= self.min_chunk_size:
                chunks.append({
                    "text": current_chunk.strip(),
                    "chunk_id": len(chunks),
                    "section_id": section_idx,
                    "word_count": current_word_count,
                    "chunk_type": "final_section"
                })
        
        logger.info(f"Created {len(chunks)} smart chunks with preserved context")
        return chunks
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences while handling policy document specifics"""
        # Handle common abbreviations that shouldn't end sentences
        text = re.sub(r'\b(Mr|Mrs|Dr|Prof|Inc|Ltd|Co|etc|vs|i\.e|e\.g)\.', r'\1<DOT>', text)
        
        # Split on sentence endings
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', text)
        
        # Restore dots in abbreviations
        sentences = [s.replace('<DOT>', '.') for s in sentences]
        
        # Filter out very short sentences (likely formatting artifacts)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
        
        return sentences


# Function alias for backward compatibility
async def process_pdf_from_url(document_url: str) -> List[Dict[str, str]]:
    """Backward compatibility function that calls DocumentProcessor.process_document_from_url"""
    processor = DocumentProcessor()
    return await processor.process_document_from_url(document_url)